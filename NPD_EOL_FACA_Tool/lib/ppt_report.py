"""PPT 报告生成(python-pptx)。

结构:
1. 封面:标题 / 生成时间 / SN 数量 / 汇总表(SN | 批號 | 線體 | 測試結果 | 站位數 | 圖片數 | 狀態)
2. 每个 SN 一页起:
   - 汇总信息(批號/包號/線體/SFC/測試結果/sensorID/flexid 等)
   - 站位轨迹表(站位 | 進站時間 | 機台號 | 載板號 | 穴位號 | 圖片數),行多时分页
   - 组件绑定表
   - PR 图片页(2x3 网格)
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .models import ImageRecord, SnRecord, StationRecord


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.4)
CONTENT_W = SLIDE_W - 2 * MARGIN

HEADER_FILL = RGBColor(0x1F, 0x4E, 0x79)
HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
ZEBRA_FILL = RGBColor(0xED, 0xF3, 0xFA)
TITLE_COLOR = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0xC0, 0x50, 0x4D)


def _set_run_font(run, font_name: str, size: float, bold: bool = False, color=None):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    # 中文字体:同时设置 east asian typeface
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font_name)


def _textbox(slide, x, y, w, h, text: str, font_name: str, size: float,
             bold: bool = False, color=None, align: PP_ALIGN = PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, font_name, size, bold, color)
    return box


def _add_title(slide, text: str, font_name: str, subtitle: str = ""):
    _textbox(slide, MARGIN, Inches(0.15), CONTENT_W, Inches(0.5),
             text, font_name, 24, bold=True, color=TITLE_COLOR)
    if subtitle:
        _textbox(slide, MARGIN + Inches(4.0), Inches(0.2), CONTENT_W - Inches(4.0), Inches(0.4),
                 subtitle, font_name, 11, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.RIGHT)


def _style_cell(cell, text: str, font_name: str, size: float,
                bold: bool = False, color=None, fill=None):
    cell.margin_left = cell.margin_right = Inches(0.04)
    cell.margin_top = cell.margin_bottom = Inches(0.01)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    _set_run_font(run, font_name, size, bold, color)


def _add_table(slide, x, y, w, headers: Sequence[str], rows: Sequence[Sequence[str]],
               font_name: str, size: float = 9.5, row_height: float = 0.28,
               col_widths: Optional[Sequence[float]] = None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(row_height * n_rows))
    table = shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(w.inches * cw / total)
    table.rows[0].height = Inches(row_height + 0.04)
    for j, h in enumerate(headers):
        _style_cell(table.cell(0, j), h, font_name, size, bold=True,
                    color=HEADER_TEXT, fill=HEADER_FILL)
    for i, row in enumerate(rows, start=1):
        table.rows[i].height = Inches(row_height)
        for j, val in enumerate(row):
            if j >= n_cols:
                break
            fill = ZEBRA_FILL if i % 2 == 0 else None
            _style_cell(table.cell(i, j), val, font_name, size, fill=fill)
    return shape


def _paginate(items: Sequence[Any], per_page: int) -> List[Sequence[Any]]:
    return [items[i : i + per_page] for i in range(0, len(items), per_page)] or [[]]


def _summary_pairs(rec: SnRecord) -> List[Tuple[str, str]]:
    pairs = []
    for key in ("批號", "線體", "包號", "SFC", "EOL測試結果", "FOL測試結果",
                "包裝時間", "箱號", "出貨地址", "出貨時間"):
        for k, v in rec.summary.items():
            if key in k and v:
                pairs.append((k, v))
                break
    if rec.sensor_id:
        pairs.append(("sensorID", rec.sensor_id))
    if rec.flex_id:
        pairs.append(("flexid", rec.flex_id))
    if not pairs:
        pairs = list(rec.summary.items())[:12]
    return pairs


def _station_table_rows(rec: SnRecord) -> List[List[str]]:
    rows = []
    for st in rec.stations:
        rows.append([
            st.station,
            st.time or "-",
            st.mc_id or "-",
            st.carrier or "-",
            st.pocket or "-",
            str(st.image_count()),
        ])
    return rows


def _component_rows(rec: SnRecord) -> List[List[str]]:
    return [
        [c.material, c.id, c.name, c.station]
        for c in rec.components
    ]


def _overview_rows(records: List[SnRecord]) -> List[List[str]]:
    rows = []
    for rec in records:
        lot = ""
        for k, v in rec.summary.items():
            if "批" in k and v:
                lot = v
                break
        line = ""
        for k, v in rec.summary.items():
            if "線體" in k or "线体" in k:
                line = v
                break
        result = ""
        for k, v in rec.summary.items():
            if "測試結果" in k or "测试结果" in k:
                result = v
                break
        rows.append([
            rec.sn,
            lot or "-",
            line or "-",
            result or "-",
            str(len(rec.stations)),
            str(len(rec.all_images())),
            "OK" if not rec.errors else f"WARN({len(rec.errors)})",
        ])
    return rows


def _station_slide(slide, rec: SnRecord, page_rows: Sequence[Sequence[str]],
                   page_no: int, total_pages: int, font_name: str, y_start: float):
    y = y_start
    _textbox(slide, MARGIN, Inches(y), CONTENT_W, Inches(0.28),
             f"站位轨迹 {page_no}/{total_pages}", font_name, 13, bold=True, color=TITLE_COLOR)
    _add_table(
        slide,
        MARGIN,
        Inches(y + 0.28),
        CONTENT_W,
        ["#", "站位", "進站時間", "機台號", "載板號", "穴位號", "圖片數"],
        [[str(i + 1 + (page_no - 1) * 13), *list(r)] for i, r in enumerate(page_rows)],
        font_name,
        size=9,
        row_height=0.26,
        col_widths=[0.5, 2.6, 2.4, 2.4, 2.0, 2.0, 1.1],
    )


class _PerSnBuilder:
    """每个 SN 的幻灯片构建(拆分出去便于组织分页)。"""

    def __init__(self, prs: Presentation, rec: SnRecord, font_name: str,
                 report_cfg: Dict[str, Any]):
        self.prs = prs
        self.rec = rec
        self.font = font_name
        self.cfg = report_cfg
        self.blank = prs.slide_layouts[6]

    def add_slide(self):
        return self.prs.slides.add_slide(self.blank)

    def build(self):
        self._summary_slide()
        self._station_slides()
        self._component_slides()
        self._image_slides()

    def _summary_slide(self):
        rec = self.rec
        slide = self.add_slide()
        _add_title(slide, f"SN: {rec.sn}", self.font,
                   subtitle=" / ".join(rec.errors) if rec.errors else "")
        pairs = _summary_pairs(rec)
        if pairs:
            rows = [[k, v] for k, v in pairs]
            _textbox(slide, MARGIN, Inches(0.7), CONTENT_W, Inches(0.28),
                     "SN 汇总信息", self.font, 13, bold=True, color=TITLE_COLOR)
            _add_table(slide, MARGIN, Inches(0.98), Inches(6.2), ["項目", "內容"],
                       rows, self.font, size=9.5, row_height=0.26,
                       col_widths=[2.4, 3.8])
        station_count = len(rec.stations)
        img_count = len(rec.all_images())
        _textbox(slide, Inches(7.0), Inches(0.7), Inches(6.0), Inches(0.28),
                 f"站位数: {station_count}    图片数: {img_count}", self.font,
                 12, bold=True, color=ACCENT)
        if rec.raw_files:
            _textbox(slide, Inches(7.0), Inches(1.05), Inches(6.0), Inches(0.5),
                     f"原始页面: {rec.raw_files[0]}", self.font, 8,
                     color=RGBColor(0x88, 0x88, 0x88))

    def _station_slides(self):
        rows = _station_table_rows(self.rec)
        pages = _paginate(rows, 13)
        total = len(pages)
        for i, page in enumerate(pages, start=1):
            slide = self.add_slide()
            _add_title(slide, f"SN: {self.rec.sn} - 站位轨迹", self.font)
            _station_slide(slide, self.rec, page, i, total, self.font, y_start=0.7)

    def _component_slides(self):
        rows = _component_rows(self.rec)
        if not rows:
            return
        for i, page in enumerate(_paginate(rows, 16), start=1):
            slide = self.add_slide()
            _add_title(slide, f"SN: {self.rec.sn} - 组件绑定", self.font)
            _textbox(slide, MARGIN, Inches(0.7), CONTENT_W, Inches(0.28),
                     f"组件绑定 {i}/{len(_paginate(rows, 16))}", self.font,
                     13, bold=True, color=TITLE_COLOR)
            _add_table(slide, MARGIN, Inches(0.98), CONTENT_W,
                       ["材料", "ID", "名稱", "使用站位"], page, self.font,
                       size=9.5, row_height=0.26,
                       col_widths=[2.5, 4.5, 3.5, 3.5])

    def _image_slides(self):
        images = self.rec.all_images()
        if not images:
            return
        pages = _paginate(images, 6)
        for i, page in enumerate(pages, start=1):
            slide = self.add_slide()
            _add_title(slide, f"SN: {self.rec.sn} - PR 图片 ({i}/{len(pages)})", self.font)
            self._image_grid(slide, page)

    def _image_grid(self, slide, images: Sequence[ImageRecord]):
        cols, rows = 3, 2
        gap_x = Inches(0.25)
        gap_y = Inches(0.2)
        top = Inches(0.75)
        cell_w = (CONTENT_W - gap_x * (cols - 1)) / cols
        cell_h = Inches(3.0)
        for idx, img in enumerate(images):
            r, c = divmod(idx, cols)
            x = MARGIN + c * (cell_w + gap_x)
            y = top + r * (cell_h + gap_y)
            if img.local_path and Path(img.local_path).exists():
                try:
                    slide.shapes.add_picture(img.local_path, x, y, width=cell_w, height=cell_h)
                except Exception:
                    _textbox(slide, x, y, cell_w, cell_h,
                             f"[图片不可用]\n{img.filename}", self.font, 9,
                             color=ACCENT)
            else:
                _textbox(slide, x, y, cell_w, cell_h,
                         f"{img.station}\n{img.filename or img.url[:60]}", self.font, 9,
                         color=RGBColor(0x66, 0x66, 0x66))
            _textbox(slide, x, y + cell_h - Inches(0.22), cell_w, Inches(0.2),
                     f"{img.station} | {img.img_type or ''}".strip(" |"), self.font, 8,
                     bold=True, color=TITLE_COLOR)


def _build_per_sn(prs: Presentation, rec: SnRecord, font_name: str,
                  report_cfg: Dict[str, Any]) -> None:
    _PerSnBuilder(prs, rec, font_name, report_cfg).build()


def build_ppt(records: List[SnRecord], out_path: Path,
              cfg: Dict[str, Any], run_info: Optional[Dict[str, Any]] = None) -> Path:
    report_cfg = cfg.get("report", {})
    font_name = report_cfg.get("slide_font", "Microsoft YaHei")
    run_info = run_info or {}

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    _textbox(cover, MARGIN, Inches(1.2), CONTENT_W, Inches(0.9),
             "SN 全制程追溯报告", font_name, 40, bold=True, color=TITLE_COLOR,
             align=PP_ALIGN.CENTER)
    _textbox(cover, MARGIN, Inches(2.1), CONTENT_W, Inches(0.5),
             f"生成时间: {datetime.now().strftime('%Y/%m/%d %H:%M:%S')}", font_name, 14,
             color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)
    _textbox(cover, MARGIN, Inches(2.6), CONTENT_W, Inches(0.5),
             f"SN 数量: {len(records)}", font_name, 14,
             color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)

    if records:
        ov = prs.slides.add_slide(blank)
        _add_title(ov, "SN 汇总", font_name)
        rows = _overview_rows(records)
        for i, page in enumerate(_paginate(rows, 28), start=1):
            if i > 1:
                ov = prs.slides.add_slide(blank)
                _add_title(ov, "SN 汇总(续)", font_name)
            _add_table(
                ov, MARGIN, Inches(0.75), CONTENT_W,
                ["SN", "批號", "線體", "測試結果", "站位數", "圖片數", "狀態"],
                page, font_name, size=9, row_height=0.26,
                col_widths=[3.4, 2.0, 1.8, 2.2, 1.2, 1.2, 1.6],
            )

    for rec in records:
        _build_per_sn(prs, rec, font_name, report_cfg)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
