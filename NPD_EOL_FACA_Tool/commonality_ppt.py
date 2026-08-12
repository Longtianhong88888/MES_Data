#!/usr/bin/env python3
"""共性分析结果汇总 PPT:桑基图 + Top20 分析 + 算法说明。

用法:
    python commonality_ppt.py --data output/commonality_NPI15.csv
        --project BOI-T --fail-count 15 --mode npi
        --out output/共性分析报告_BOI-T_NPI15.pptx
"""

from __future__ import annotations

import argparse
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import pandas as pd

from commonality_analysis import (
    analyze_commonality, apriori_rules, pick_factor_columns, _clean,
    chart_category, CATEGORY_NAMES, CATEGORY_ORDER,
    CHART_MACHINE, CHART_MATERIAL, CHART_HEAD, CHART_TIME,
    CHART_WAFER, CHART_CARRIER,
)


# ---------- 桑基图(matplotlib 手绘,无第三方 sankey 依赖) ----------
def draw_sankey(
    dimension: str,
    df: pd.DataFrame,
    top_values: Optional[List[str]],
    out_png: Path,
    pass_scale: float = 0.04,
    max_values: int = 7,
) -> str:
    """Fail/Pass 流向桑基图(对齐工具源码语义):
    - 节点按 z 值显著性着色:z<1.96 灰;≥1.96 用 5 档红绿渐变
      (绿→黄→橙→红,按权重分箱);fail_ratio>0.9/0.95/0.99 逐级加深
    - 流线:Fail=salmon,Pass=silver
    - 节点标注: 值(Ratio:xx%) / FR:xx%(fail/total)
    z = sqrt(n) × (该类Fail率-整体Fail率) / sqrt(整体率×(1-整体率))
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib as mpl
    mpl.rcParams["font.sans-serif"] = [
        "Microsoft YaHei", "PingFang SC", "Arial Unicode MS",
        "SimHei", "DejaVu Sans"]
    mpl.rcParams["axes.unicode_minus"] = False
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyBboxPatch, PathPatch
    from matplotlib.path import Path as MPath
    import bisect
    import math as _math

    is_fail = df["pass_fail"].astype(str).str.strip().isin(("fail", "Fail"))
    vals = df[dimension].map(_clean)
    total_fail = int(is_fail.sum())
    total_pass = int((~is_fail).sum())
    overall_fr = total_fail / len(df) if len(df) else 0.0

    def zscore(n_total, n_fail):
        if n_total <= 0 or overall_fr <= 0 or overall_fr >= 1:
            return 0.0
        fr = n_fail / n_total
        return (_math.sqrt(n_total) * (fr - overall_fr)
                / _math.sqrt(overall_fr * (1 - overall_fr)))

    def node_color(z, fail_ratio, z_max):
        fixed = ["#24F20C", "#B5F20C", "#FFD700", "#FF8C00", "#FF3333"]
        if z < 1.96:
            return "lightgrey"
        if fail_ratio > 0.99:
            z = z_max + 8
        elif fail_ratio >= 0.95:
            z = z_max + 5
        elif fail_ratio >= 0.9:
            z = z_max + 3
        weights = [3, 2.5, 2, 1.5, 1]
        total_w = sum(weights)
        cum = [1.96 + (z_max - 1.96) * sum(weights[:k]) / total_w
               for k in range(1, 6)]
        idx = bisect.bisect_left(cum, z)
        return fixed[max(0, min(idx, 4))]

    if top_values:
        shown = top_values[:max_values]
    else:
        # 未指定值列表:取该维度全部值,按共性(z 显著性)从高到低排,
        # 机台等全部串在一张图(与参考工具"共性从高到低"一致)
        vc = vals[is_fail].value_counts()
        cand = [v for v in vc.index if str(v) not in ("", "nan")]

        def _z(v):
            m = vals == v
            nf = int((is_fail & m).sum())
            return zscore(int(m.sum()), nf)

        shown = sorted(cand, key=_z, reverse=True)
    items = []  # (label, fail, pass)
    for v in shown:
        m = vals == v
        items.append((v, int((is_fail & m).sum()), int((~is_fail & m).sum())))
    other_fail = total_fail - sum(i[1] for i in items)
    other_pass = total_pass - sum(i[2] for i in items)
    if top_values and (other_fail > 0 or other_pass > 0):
        # 仅指定值列表(截断展示)时聚合"其他"
        items.append(("其他", other_fail, other_pass))
    # 计算各值 z 与颜色
    zmax = max([zscore(f + p, f) for _, f, p in items] or [0.0])
    item_meta = []
    for label, nf, np_ in items:
        n = nf + np_
        z = zscore(n, nf)
        ratio = nf / total_fail if total_fail else 0.0
        item_meta.append((label, nf, np_, z, ratio, node_color(z, ratio, zmax)))

    # 布局
    n_mid = len(items)
    x0, x1, x2 = 0.13, 0.5, 0.87
    y_mid_top = 0.92
    mid_h = 0.82 / max(n_mid, 1)
    fig, ax = plt.subplots(figsize=(12, 0.62 * n_mid + 2.2))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    C_FAIL = "#E74C3C"
    C_PASS = "#2ECC71"
    LINK_FAIL = "#FA8072"   # salmon
    LINK_PASS = "#C0C0C0"   # silver

    def node(x, y, h, w, label, fc, ec="white"):
        box = FancyBboxPatch((x - w / 2, y), w, h,
                             boxstyle="round,pad=0.008",
                             fc=fc, ec=ec, lw=0.8)
        ax.add_patch(box)
        ax.text(x, y + h / 2, label, ha="center", va="center",
                fontsize=9, color="white", fontweight="bold")

    def flow(xa, ya, xb, yb, width, color, alpha=0.55):
        w = max(width, 0.0015)
        mid = (xa + xb) / 2
        # 上边缘(起→止) + 下边缘(止→起) 闭合带
        verts = [
            (xa, ya + w / 2), (mid, ya + w / 2), (mid, yb + w / 2),
            (xb, yb + w / 2), (xb, yb - w / 2), (mid, yb - w / 2),
            (mid, ya - w / 2), (xa, ya - w / 2),
        ]
        codes = [MPath.MOVETO, MPath.CURVE4, MPath.CURVE4, MPath.CURVE4,
                 MPath.LINETO, MPath.CURVE4, MPath.CURVE4, MPath.CURVE4]
        path = PathPatch(MPath(verts, codes), lw=0, fc=color,
                         alpha=alpha)
        ax.add_patch(path)

    # 左节点:总 Fail / 总 Pass
    node(x0, 0.86, 0.1, 0.16, "Fail %d" % total_fail, C_FAIL)
    node(x0, 0.12, 0.1, 0.16, "Pass %d" % total_pass, C_PASS)

    # 中间节点:因子值(按工具配色:z 显著性)
    mid_centers = []
    for i, (label, nf, np_, z, ratio, fc) in enumerate(item_meta):
        y = y_mid_top - (i + 1) * mid_h
        label_txt = "%s\nRatio:%.1f%%  FR:%.1f%%(%d/%d)" % (
            label, ratio * 100,
            (nf / n if n else 0) * 100, nf, n)
        node(x1, y, mid_h * 0.9, 0.36, label_txt, fc)
        mid_centers.append((y + mid_h * 0.45, nf, np_))

    # 右节点:Fail / Pass
    node(x2, 0.82, 0.12, 0.14, "Fail", C_FAIL)
    node(x2, 0.08, 0.12, 0.14, "Pass", C_PASS)

    # 左→中 流
    for i, (y, nf, np_) in enumerate(mid_centers):
        flow(x0, 0.91, x1, y, nf / total_fail * 0.5, LINK_FAIL, 0.5)
        flow(x0, 0.17, x1, y - mid_h * 0.4, np_ / total_pass * 0.5 * pass_scale,
             LINK_PASS, 0.35)
    # 中→右 流
    for i, (y, nf, np_) in enumerate(mid_centers):
        flow(x1, y, x2, 0.88, nf / total_fail * 0.5, LINK_FAIL, 0.5)
        flow(x1, y - mid_h * 0.4, x2, 0.14,
             np_ / total_pass * 0.5 * pass_scale, LINK_PASS, 0.35)

    ax.text(0.5, 0.975, "共性因子: %s(Fail 流向, 节点按 z 值显著性着色)" % dimension,
            ha="center", fontsize=12, fontweight="bold")
    ax.text(0.5, 0.018,
            "节点颜色: 灰=z<1.96(不显著); 绿→黄→橙→红=z 显著性递增; "
            "Fail 流 salmon / Pass 流 silver; Pass 流按 1/%.0f 缩放" % (1 / pass_scale),
            ha="center", fontsize=8, color="#7F8C8D")
    plt.tight_layout()
    fig.savefig(out_png, dpi=150)
    plt.close(fig)
    return str(out_png)


# ---------- 6 类共性图表(对齐参考工具:机台/材料/Head/时间/Wafer/穴位) ----------
CHART_NAMES = {
    CHART_MACHINE: "机台共性(桑基图)",
    CHART_MATERIAL: "材料共性(横条图)",
    CHART_HEAD: "Head 共性(柱状图)",
    CHART_TIME: "时间共性(面积图)",
    CHART_WAFER: "Wafer 共性(Wafer Map)",
    CHART_CARRIER: "穴位共性(HeatMap)",
}


def _mpl_setup():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib as mpl
    mpl.rcParams["font.sans-serif"] = [
        "Microsoft YaHei", "PingFang SC", "Arial Unicode MS",
        "SimHei", "DejaVu Sans"]
    mpl.rcParams["axes.unicode_minus"] = False


def _clean(v) -> str:
    s = str(v).strip()
    return "" if s in ("{}", "nan", "None", "") else s


def classify_chart_cols(df) -> Dict[str, List[str]]:
    """把合并表列名归到 6 类图表(按命名规则)。"""
    out: Dict[str, List[str]] = {k: [] for k in CHART_NAMES}
    for c in df.columns:
        cat = chart_category(c)
        if cat in out:
            out[cat].append(c)
    return out


def _category_stats(df, col) -> pd.DataFrame:
    """某列各取值的 投入/fail/fail率/z 统计。"""
    is_fail = df["pass_fail"].astype(str).str.strip().isin(("fail", "Fail"))
    all_fail = int(is_fail.sum())
    all_rate = all_fail / len(df) if len(df) else 0.0
    vals = df[col].map(_clean)
    rows = []
    for v, g in df.groupby(vals):
        if not v:
            continue
        n = len(g)
        nf = int(is_fail[g.index].sum())
        fr = nf / n if n else 0.0
        z = (n ** 0.5 * (fr - all_rate)
             / (all_rate * (1 - all_rate)) ** 0.5
             if 0 < all_rate < 1 else 0.0)
        rows.append({"value": v, "total": n, "fail": nf,
                     "fail_rate": fr,
                     "ratio": nf / all_fail if all_fail else 0.0,
                     "z": z})
    d = pd.DataFrame(rows)
    if not d.empty:
        d = d.sort_values("z", ascending=False).reset_index(drop=True)
    return d


def _z_color(z, ratio, zmax=None):
    """与桑基图一致的红绿渐变(z<1.96 灰,否则绿→黄→橙→红)。"""
    fixed = ["#24F20C", "#B5F20C", "#FFD700", "#FF8C00", "#FF3333"]
    if z < 1.96:
        return "lightgrey"
    if zmax is None:
        zmax = max(z, 1.96)
    weights = [3, 2.5, 2, 1.5, 1]
    cum = [1.96 + (zmax - 1.96) * sum(weights[:k]) / sum(weights)
           for k in range(1, 6)]
    idx = max(0, min(sum(1 for c in cum if z >= c), 4))
    return fixed[idx]


def draw_material_barh(df, dims, out_png) -> str:
    """材料共性:每个材料维度一行,Top5 取值按共性着色横排(对齐参考工具)。"""
    import matplotlib.pyplot as plt
    from matplotlib.patches import Patch
    _mpl_setup()
    items = []
    for dim in dims:
        st = _category_stats(df, dim)
        if st.empty:
            continue
        items.append((dim, st.head(5)))
    if not items:
        return ""
    fig, ax = plt.subplots(figsize=(14, 0.9 * len(items) + 2.5))
    zmax = max(float(t["z"].max()) for _, t in items if len(t))
    for i, (dim, top) in enumerate(items):
        item_total = int(top["total"].sum())
        for j, (_, r) in enumerate(top.iterrows()):
            ax.barh(i, 1, height=0.7, left=j, color=_z_color(
                float(r["z"]), float(r["ratio"]), zmax),
                edgecolor="black", linewidth=0.4)
            label = "%s %s %s(%d/%d)" % (
                str(r["value"]).split(".0")[0],
                "Ratio:%.1f%%" % (float(r["ratio"]) * 100),
                "FR:%.1f%%" % (float(r["fail_rate"]) * 100),
                int(r["fail"]), int(r["total"]))
            ax.text(j + 0.5, i, label, ha="center", va="center",
                    fontsize=8, color="black")
        ax.text(-0.3, i, "%s(共%d)" % (dim, item_total), ha="right",
                va="center", fontsize=9, fontweight="bold")
    ax.set_yticks(range(len(items)))
    ax.set_yticklabels([""] * len(items))
    ax.set_xticks([])
    ax.set_xlim(-0.35, 5.3)
    ax.set_ylim(-0.6, len(items) - 0.4)
    ax.set_title("Material_Commonality(材料共性,每维度取共性最高前5)",
                 fontsize=13, fontweight="bold")
    ax.legend(handles=[
        Patch(fc="#FF3333", label="共性高(红)"),
        Patch(fc="#FFD700", label="中(黄)"),
        Patch(fc="#24F20C", label="低(绿)"),
        Patch(fc="lightgrey", label="z<1.96 不显著")],
        loc="upper right", fontsize=8)
    plt.tight_layout()
    fig.savefig(str(out_png), dpi=150, bbox_inches="tight")
    plt.close(fig)
    return str(out_png)


def draw_head_barchart(df, dim, out_png, top_n=12) -> str:
    """Head 共性:各 Head_ID 的投入(蓝)/fail(橙)柱状图(对齐参考第17页)。"""
    import matplotlib.pyplot as plt
    _mpl_setup()
    st = _category_stats(df, dim)
    if st.empty:
        return ""
    st = st.head(top_n).sort_values("total", ascending=True)
    labels = [str(v)[:14] for v in st["value"]]
    x = list(range(len(st)))
    fig, ax = plt.subplots(figsize=(13, max(3.2, 0.42 * len(st) + 1.5)))
    ax.barh(x, st["total"], height=0.6, color="#3B82F6", label="投入数量")
    ax.barh(x, st["fail"], height=0.6, color="#F97316", label="Fail 数量")
    for i, (_, r) in enumerate(st.iterrows()):
        ax.text(int(r["total"]) + 2, i, "%d" % int(r["total"]),
                va="center", fontsize=8, color="#2563EB")
        if int(r["fail"]) > 0:
            ax.text(int(r["fail"]) - 2, i, "%d" % int(r["fail"]),
                    va="center", ha="right", fontsize=8, color="#C2410C")
            ax.text(int(r["total"]) * 1.02, i,
                    "FR:%.2f%% 占不良:%.1f%%" % (
                        float(r["fail_rate"]) * 100,
                        float(r["ratio"]) * 100),
                    va="center", fontsize=8, color="#6B7280")
    ax.set_yticks(x)
    ax.set_yticklabels(labels, fontsize=8)
    ax.set_xlabel("数量", fontsize=10)
    ax.set_title("Head_ID 共性: %s(蓝=投入,橙=Fail)" % dim,
                 fontsize=13, fontweight="bold")
    ax.legend(loc="lower right", fontsize=9)
    ax.xaxis.set_major_locator(plt.MaxNLocator(integer=True))
    plt.tight_layout()
    fig.savefig(str(out_png), dpi=150, bbox_inches="tight")
    plt.close(fig)
    return str(out_png)


def draw_time_chart(df, col, out_png) -> str:
    """时间共性:生产投入面积图 + Fail 红点堆叠(对齐参考第18页)。"""
    import matplotlib.pyplot as plt
    import pandas as pd
    from matplotlib.dates import DateFormatter
    _mpl_setup()
    sub = df[["pass_fail", col]].copy()
    sub = sub[sub[col].notna() & (sub[col] != "{}")]
    sub[col] = pd.to_datetime(sub[col], errors="coerce")
    sub = sub.dropna(subset=[col]).sort_values(col)
    if sub.empty:
        return ""
    span_h = (sub[col].max() - sub[col].min()).total_seconds() / 3600
    if span_h <= 24:
        freq, fail_freq = "10min", "10min"
    elif span_h <= 120:
        freq, fail_freq = "12h", "30min"
    elif span_h <= 480:
        freq, fail_freq = "2D", "30min"
    else:
        freq, fail_freq = "3D", "30min"
    sampled = sub.set_index(col).resample(freq).size()
    fig, ax = plt.subplots(figsize=(15, 6))
    ax.fill_between(sampled.index, sampled.values, alpha=0.5,
                    color="paleturquoise")
    ax.plot(sampled.index, sampled.values, linewidth=2,
            color="paleturquoise", alpha=0.8)
    ax.set_title("%s(时间共性:投入面积 + Fail 红点)" % col,
                 fontsize=13, fontweight="bold")
    ax.set_xlabel("Time", fontsize=11)
    ax.set_ylabel("投入数量", fontsize=11)
    ax.yaxis.set_major_locator(plt.MaxNLocator(integer=True))
    ax.set_ylim(bottom=0)
    ax.xaxis.set_major_formatter(DateFormatter("%m-%d %H:%M"))
    plt.setp(ax.xaxis.get_majorticklabels(), rotation=45, ha="right")
    fail_df = sub[sub["pass_fail"].astype(str).str.strip().isin(
        ("fail", "Fail"))].copy()
    if len(fail_df):
        fail_df["tg"] = fail_df[col].dt.floor(fail_freq)
        fail_df["stack"] = fail_df.groupby("tg").cumcount() + 1
        max_bin = fail_df.groupby("tg").size().max() or 1
        ax2 = ax.twinx()
        ax2.set_ylim(0, max_bin * 15)
        ax2.set_ylabel("Fail 数量(堆叠)", fontsize=11)
        ax2.scatter(fail_df[col], fail_df["stack"], marker="o",
                    s=6, alpha=0.6, c="red")
    plt.tight_layout()
    fig.savefig(str(out_png), dpi=150, bbox_inches="tight")
    plt.close(fig)
    return str(out_png)


def _parse_xy(v) -> Optional[Tuple[int, int]]:
    s = _clean(v)
    if not s or "_" not in s:
        return None
    try:
        a, b = s.split("_", 1)
        return int(float(a)), int(float(b))
    except ValueError:
        return None


def draw_wafer_map(df, xy_col, out_png) -> str:
    """Wafer 共性:Wafer Map,绿=Pass 位,红绿渐变=Fail 数(对齐参考工具)。"""
    import matplotlib.pyplot as plt
    _mpl_setup()
    is_fail = df["pass_fail"].astype(str).str.strip().isin(("fail", "Fail"))
    pos = df[xy_col].map(_parse_xy)
    sub = df.loc[pos.notna()].copy()
    sub["_xy"] = pos[pos.notna()]
    sub["_fail"] = is_fail[pos.notna()]
    if sub.empty:
        return ""
    groups = sub.groupby("_xy")
    fail_counts = {k: int(g["_fail"].sum()) for k, g in groups}
    max_fail = max(fail_counts.values()) or 1
    fig, ax = plt.subplots(figsize=(12, 8))
    pass_only = [k for k, n in fail_counts.items() if n == 0]
    if pass_only:
        px, py = zip(*pass_only)
        ax.scatter(px, py, c="green", marker="s", alpha=0.6, s=30,
                   label="Pass 位")
    for (x, y), n in fail_counts.items():
        if n == 0:
            continue
        norm = min(n / max_fail, 1.0)
        if norm <= 0.3:
            g = 1 - 0.3 * (norm / 0.3)
        elif norm <= 0.7:
            g = 0.7 - 0.4 * ((norm - 0.3) / 0.4)
        else:
            g = 0.3 - 0.3 * ((norm - 0.7) / 0.3)
        g = max(0.0, min(1.0, g))
        ax.scatter(x, y, c=[(1, g, 0)], marker="s", alpha=0.8,
                   s=65 if n == max_fail else 47)
    ax.set_title("Wafer_Map(%s)" % xy_col, fontsize=13, fontweight="bold")
    ax.set_xlabel(xy_col.split("_")[0] + "_X", fontsize=11)
    ax.set_ylabel(xy_col.split("_")[0] + "_Y", fontsize=11)
    ax.legend(loc="upper right", fontsize=9)
    plt.tight_layout()
    fig.savefig(str(out_png), dpi=150, bbox_inches="tight")
    plt.close(fig)
    return str(out_png)


def draw_carrier_heatmap(df, col, out_png) -> str:
    """穴位共性:Carrier 穴位(row_col)Fail 率 HeatMap(共性越高越蓝)。"""
    import matplotlib.pyplot as plt
    import numpy as np
    from matplotlib.colors import LinearSegmentedColormap
    _mpl_setup()
    is_fail = df["pass_fail"].astype(str).str.strip().isin(("fail", "Fail"))
    pos = df[col].map(_parse_xy)
    sub = df.loc[pos.notna()].copy()
    sub["_row"] = [p[0] for p in pos[pos.notna()]]
    sub["_col"] = [p[1] for p in pos[pos.notna()]]
    sub["_fail"] = is_fail[pos.notna()]
    if sub.empty:
        return ""
    rate = sub.groupby(["_row", "_col"])["_fail"].mean()
    rows = sorted({p[0] for p in rate.index})
    cols = sorted({p[1] for p in rate.index})
    if not rows or not cols:
        return ""
    mat = np.full((len(rows), len(cols)), np.nan)
    for (r, c), v in rate.items():
        mat[rows.index(r)][cols.index(c)] = v
    cmap = LinearSegmentedColormap.from_list(
        "blue_heat", ["#E8F1FF", "#3B82F6", "#1E3A8A"])
    fig, ax = plt.subplots(figsize=(11, 8))
    im = ax.imshow(mat, cmap=cmap, aspect="auto", origin="lower",
                   vmin=0, vmax=max(float(rate.max()), 1e-9))
    ax.set_xticks(range(len(cols)))
    ax.set_xticklabels(cols, fontsize=8)
    ax.set_yticks(range(len(rows)))
    ax.set_yticklabels(rows, fontsize=8)
    ax.set_xlabel("穴位 X", fontsize=11)
    ax.set_ylabel("穴位 Y", fontsize=11)
    ax.set_title("Carrier 穴位共性 HeatMap: %s(越蓝共性越高)" % col,
                 fontsize=13, fontweight="bold")
    fig.colorbar(im, ax=ax, label="Fail 率")
    plt.tight_layout()
    fig.savefig(str(out_png), dpi=150, bbox_inches="tight")
    plt.close(fig)
    return str(out_png)


# ---------- PPT ----------
def _add_title(slide, text, size=30):
    box = slide.shapes.add_textbox(360000, 120000, 12192000, 800000)
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = __import__("pptx").util.Pt(size)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = __import__("pptx").dml.color.RGBColor(
        0x1F, 0x3A, 0x5F)


def _add_bullets(slide, lines, top=900000, size=15):
    box = slide.shapes.add_textbox(500000, top, 12300000, 5200000)
    tf = box.text_frame
    tf.word_wrap = True
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.font.size = Pt(size)
        p.space_after = Pt(6)
    return box


def _add_table(slide, headers, rows, left=360000, top=900000,
               width=12600000, height=5600000):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                   left, top, width, height)
    tbl = shape.table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
        c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        c.text_frame.paragraphs[0].font.size = Pt(11)
        c.text_frame.paragraphs[0].font.bold = True
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(v)
            c.text_frame.paragraphs[0].font.size = Pt(10)
            if j == 0 and str(row[-1]) == "重点":
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xFF, 0xC7, 0xCE)
    return shape


def _pick_chart_dims(top_df: pd.DataFrame, df: pd.DataFrame
                     ) -> Dict[str, Optional[str]]:
    """为 6 类图表各选一个代表列:优先 top20 里的,否则取该类型第一列。"""
    chart_cols = classify_chart_cols(df)
    top_dims = [str(d) for d in top_df["dimension"]] if len(top_df) else []
    n = max(len(df), 1)

    def parse_ratio(col, parser):
        try:
            vals = df[col].dropna()
            if len(vals) == 0:
                return 0.0
            return sum(1 for v in vals.head(200) if parser(v) is not None) / min(len(vals), 200)
        except Exception:
            return 0.0

    picks: Dict[str, Optional[str]] = {}
    for kind in CHART_NAMES:
        cols = chart_cols.get(kind, [])
        col = None
        for c in cols:
            if c in top_dims:
                col = c
                break
        if col is None and cols:
            if kind == CHART_CARRIER:
                col = max(cols, key=lambda c: parse_ratio(c, _parse_xy))
            elif kind == CHART_WAFER:
                col = max(cols, key=lambda c: parse_ratio(c, _parse_xy))
            elif kind == CHART_TIME:
                def _dt(v):
                    try:
                        import pandas as pd
                        return pd.to_datetime(v, errors="coerce")
                    except Exception:
                        return None
                col = max(cols, key=lambda c: parse_ratio(c, _dt))
            elif kind == CHART_MATERIAL:
                col = max(cols, key=lambda c: len(_category_stats(df, c)))
            else:
                col = cols[0]
        picks[kind] = col
    return picks


def _add_picture_fit(slide, png: Path, top=1.35, box_w=11.3, box_h=5.7):
    """图片按幻灯片可用区域等比缩放,避免超高图溢出。"""
    from pptx.util import Inches
    from PIL import Image as _Img
    with _Img.open(png) as _im:
        _iw, _ih = _im.size
    _scale = min(box_w / _iw, box_h / _ih)
    slide.shapes.add_picture(str(png), Inches(1.0), Inches(top),
                             width=Inches(_iw * _scale),
                             height=Inches(_ih * _scale))


def build_ppt(project: str, mode: str, fail_count: int,
              top_df: pd.DataFrame, rules_df: pd.DataFrame,
              df: pd.DataFrame, out: Path,
              sankey_dims: Optional[List[str]] = None) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    ts = datetime.now().strftime("%Y-%m-%d")

    # 1) 封面
    s = prs.slides.add_slide(blank)
    _add_title(s, "全制程共性分析报告", 40)
    _add_bullets(s, [
        "专案: %s    分析模式: %s" % (project, mode.upper()),
        "Fail SN 数量: %d    生成日期: %s" % (fail_count, ts),
        "数据源: Load_DataCenterData(全量 + Pass 采样基线)",
    ], top=2600000)

    # 2) 算法说明
    s = prs.slides.add_slide(blank)
    _add_title(s, "Top20 共性评分算法说明")
    _add_bullets(s, [
        "① 因子化: 每个维度(机台/料批/供应商/治工具/工作区…)的每个值,统计其 Fail 数与样本数",
        "② Lift = 该值Fail率 ÷ 整体Fail率,衡量“异常放大”倍数; Fail占比 = 该值Fail数 ÷ 总Fail数",
        "③ Score = Fail占比 × log2(Lift),越高越可疑;覆盖>80%且Lift低的项自动降权(必要非充分条件)",
        "④ 显著性: Fisher精确检验(p值),与Pass基线对比; MP模式再做FDR多重校正,NPI小样本用原始p",
        "⑤ 组合规则(第二层): Apriori挖掘Fail中共现的“维度=值”组合,按Lift排序;滤除机器×机器结构性共现",
        "⑥ 结论定位: 优先看 Score 高 + Fail覆盖大 + p值小的项;NPI阶段重点看材料/治工具/机台内部工作区",
    ], size=15)

    # 3) Top20 按 6 大类分组展示
    cols = ["维度", "值", "Fail数", "样本数", "Fail率", "Lift", "Fail占比", "p_adj", "Score"]

    def _top_rows(rows_list):
        rows = []
        for r in rows_list:
            rows.append([
                str(r.get("dimension", "")),
                str(r.get("value", "")),
                int(r.get("fail_count", 0)),
                int(r.get("unit_count", 0)),
                "%.1f%%" % (float(r.get("fail_rate", 0)) * 100),
                "%.1fx" % float(r.get("lift", 0)),
                "%.1f%%" % (float(r.get("fail_ratio", 0)) * 100),
                "%.2e" % float(r.get("p_adj", 1)),
                "%.2f" % float(r.get("score", 0)),
            ])
        return rows

    cat_groups = {k: [] for k in CATEGORY_ORDER}
    other_rows = []
    for _, r in top_df.iterrows():
        cat = chart_category(str(r.get("dimension", "")))
        if cat in cat_groups:
            cat_groups[cat].append(r)
        else:
            other_rows.append(r)

    # 3-1) 总览页
    s = prs.slides.add_slide(blank)
    _add_title(s, "Top20 共性可疑点总览(按6大类)")
    lines = ["· 共 %d 项,按 机台/材料/Head/时间/Wafer/穴位 六大类分组" % len(top_df)]
    for k in CATEGORY_ORDER:
        n = len(cat_groups[k])
        if n:
            lines.append("· %s: %d 项" % (CATEGORY_NAMES[k], n))
    if other_rows:
        lines.append("· 其他: %d 项" % len(other_rows))
    _add_bullets(s, lines, size=16)

    # 3-2) 每大类一页(超过 12 项自动分页)
    for k in CATEGORY_ORDER:
        rows = _top_rows(cat_groups[k])
        if not rows:
            continue
        for start in range(0, len(rows), 12):
            s = prs.slides.add_slide(blank)
            _add_title(s, "Top20 · %s 共性(%d项)" % (
                CATEGORY_NAMES[k], len(rows)), 26)
            _add_table(s, cols, rows[start:start + 12], height=5600000)
    if other_rows:
        rows = _top_rows(other_rows)
        for start in range(0, len(rows), 12):
            s = prs.slides.add_slide(blank)
            _add_title(s, "Top20 · 其他(%d项)" % len(rows), 26)
            _add_table(s, cols, rows[start:start + 12], height=5600000)

    # 4) 六类共性图表(机台/材料/Head/时间/Wafer/穴位,全部出现)
    picks = _pick_chart_dims(top_df, df)
    # 材料维度取共性最高前 10 个(每个维度 Top5 取值)
    material_dims = []
    for c in classify_chart_cols(df).get(CHART_MATERIAL, []):
        st = _category_stats(df, c)
        if len(st) >= 2:
            material_dims.append((float(st["z"].max()), c))
    material_dims = [c for _, c in
                     sorted(material_dims, key=lambda x: x[0], reverse=True)][:10]
    for kind in (CHART_MACHINE, CHART_MATERIAL, CHART_HEAD,
                 CHART_TIME, CHART_WAFER, CHART_CARRIER):
        col = picks.get(kind)
        if not col:
            continue
        png = Path(out).parent / ("chart_%s.png" % kind)
        if kind == CHART_MACHINE:
            if not draw_sankey(col, df, None, png):
                continue
        elif kind == CHART_MATERIAL:
            if not draw_material_barh(df, material_dims or [col], png):
                continue
        elif kind == CHART_HEAD:
            if not draw_head_barchart(df, col, png):
                continue
        elif kind == CHART_TIME:
            if not draw_time_chart(df, col, png):
                continue
        elif kind == CHART_WAFER:
            if not draw_wafer_map(df, col, png):
                continue
        else:
            if not draw_carrier_heatmap(df, col, png):
                continue
        s = prs.slides.add_slide(blank)
        title = "%s(%s)" % (CHART_NAMES[kind], str(col)[:38])
        _add_title(s, title, 22)
        _add_picture_fit(s, png)

    # 5) 附加桑基图(用户勾选维度;机台标准图已含,避免重复)
    if sankey_dims:
        std_machine = picks.get(CHART_MACHINE)
        for dim in sankey_dims:
            if dim == std_machine:
                continue
            png = Path(out).parent / ("sankey_%s.png" % dim.replace("/", "_"))
            if not draw_sankey(dim, df, None, png):
                continue
            s = prs.slides.add_slide(blank)
            _add_title(s, "桑基图: %s" % str(dim)[:40], 24)
            _add_picture_fit(s, png)

    # 6) 组合规则
    if not rules_df.empty:
        s = prs.slides.add_slide(blank)
        _add_title(s, "组合规则(第二层,已滤除机器×机器)")
        rows = []
        for _, r in rules_df.head(15).iterrows():
            rows.append([
                str(r.get("item_a", "")), str(r.get("item_b", "")),
                int(r.get("count_ab", 0)),
                "%.1f%%" % (float(r.get("support", 0)) * 100),
                "%.2fx" % float(r.get("lift", 0)),
            ])
        _add_table(s, ["组合A", "组合B", "共现Fail", "Support", "Lift"],
                   rows, height=5000000)

    # 7) 结论
    s = prs.slides.add_slide(blank)
    _add_title(s, "结论与建议")
    tops = []
    for _, r in top_df.head(5).iterrows():
        tops.append("· %s = %s (Fail %d/%d, Lift %.1fx, Score %.2f)" % (
            r.get("dimension"), r.get("value"), int(r.get("fail_count", 0)),
            fail_count, float(r.get("lift", 0)), float(r.get("score", 0))))
    _add_bullets(s, ["【重点关注】"] + tops + [
        "",
        "【建议】",
        "· 优先排查 Score 前 5 项对应的机台/料批/治工具/工作区",
        "· 组合规则中的跨维度组合(如料批×供应商)指向物料或工艺交互问题",
        "· 小样本(NPI)结论当线索,建议结合机台日志与物料追溯人工复核",
    ], size=14)

    prs.save(str(out))
    return str(out)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True, help="合并 CSV(含 pass_fail)")
    ap.add_argument("--project", default="")
    ap.add_argument("--mode", default="auto")
    ap.add_argument("--fail-count", type=int, default=0)
    ap.add_argument("--sankey-dims", default="",
                    help="桑基图维度,逗号分隔(如 FLEX_lot_ID_1,MC_ID);"
                         "留空=自动取前3个")
    ap.add_argument("--out", default="output/共性分析报告.pptx")
    args = ap.parse_args()

    df = pd.read_csv(args.data, low_memory=False, dtype=str)
    fail_count = args.fail_count or int(
        (df["pass_fail"].astype(str).str.lower() == "fail").sum())
    mode = args.mode
    if mode == "auto":
        mode = "npi" if fail_count < 20 else "mp"
    rows, _ = analyze_commonality(
        df, min_fail=1 if mode == "npi" else 3,
        fail_values=("fail", "Fail"), mode=mode)
    top = pd.DataFrame(rows).head(20)
    rules = apriori_rules(df, fail_col="pass_fail")
    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    dims = [d.strip() for d in args.sankey_dims.split(",") if d.strip()] or None
    build_ppt(args.project or "未知专案", mode, fail_count, top, rules,
              df, out, sankey_dims=dims)
    print("PPT 已生成: %s" % out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
